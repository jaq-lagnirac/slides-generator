# Justin Caringal
# Takes a directory of images and converts them into a
# .pptx presentation/slides deck

# libraries
import os
import sys
import logging
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# logging formatting
SCRIPT_PATH = os.path.abspath(__file__)
FORMAT = '[%(asctime)s] %(levelname)s %(message)s'
l = logging.getLogger()
lh = logging.StreamHandler()
lh.setFormatter(logging.Formatter(FORMAT))
l.addHandler(lh)
l.setLevel(logging.INFO)
debug = l.debug; info = l.info; warning = l.warning; error = l.error

# argparse constants
DESCRIPTION = '''
'''
EPILOG = '''
'''

# constants
ACCEPTABLE_IMG_TYPES = ['.jpg', '.png', '.jpeg']


def check_dir(dir : str) -> str:
    """checks the existence of the dir
    
    A function to ensure the directory exists before
    executing the rest of the program

    Args:
        dir (str): a relative path to the input dir
    
    Returns:
        str: Returns name of directory

    """

    if not os.path.isdir(dir):
        error(f'Input {dir} is not a directory.')
        sys.exit(1)

    if not os.path.exists(dir):
        error(f'Directory {dir} not found.')
        sys.exit(1)

    info(f'{args.dir} is a valid directory')

    # adds safety slash '/' to end of dir str in order to prevent unaccounted outputs
    if dir[-1] != '/': # if last char is not a '/', add a /
        dir += '/'
        debug('Added safety slash')
    
    dirname = os.path.split(os.path.dirname(dir))[-1] # takes last value

    debug(f'dirname: {dirname}')
    return dirname


def is_image(file : str) -> bool | str:
    """checks if file is an acceptable image format
    
    A function to check the validity of the image

    Args:
        file (str) : the basename of the file

    Returns:
        bool: Returns True if filetype is acceptable, false otherwise

    """

    root, ext = os.path.splitext(file)

    if ext.lower() not in ACCEPTABLE_IMG_TYPES:
        warning(f'{file} is not an acceptable image type.')
        debug(f'Root: {root} - Ext: {ext}')
        return False
    
    return True


def main() -> None:
    """MAIN FUNCTION"""

    # checks existence of dir, extracts dirname
    info(f'Checking input directory {args.dir}')
    dirname = check_dir(args.dir)

    # inits presentation
    info('Initializing presentation')
    prez = Presentation()
    # sets aspect ratio to standard Google Slides size
    prez.slide_width = Inches(10)
    prez.slide_height = Inches(5.625)
    
    # adds title slide
    info('Creating title slide')
    slide = prez.slides.add_slide(prez.slide_layouts[6]) # pos=6, Blank
    # generates title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)
    title = slide.shapes.add_textbox(left, top, width, height)
    tf = title.text_frame
    p = tf.add_paragraph()
    p.text = dirname
    debug(f'Title slide text: {dirname}')
    p.font.bold = True
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER
    # generates subtitle
    top = Inches(3)
    title = slide.shapes.add_textbox(left, top, width, height)
    tf = title.text_frame
    p = tf.add_paragraph()
    subtitle_text = 'Created using tools developed by Justin Caringal'
    p.text = subtitle_text
    debug(f'Subtitle text: {subtitle_text}')
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    # iterates through image directory to create slides
    left = Inches(0.5)
    top = Inches(0)
    width = Inches(8)
    height = Inches(1)
    img_top = Inches(1.1)
    img_height = Inches(4.275)
    for img_basename in os.listdir(args.dir):
        
        info(f'Creating slide for {img_basename}')

        # creates relative path to image
        img_path = os.path.join(args.dir, img_basename)
        
        # checks validity of image, skips if not valid
        if not is_image(img_path):
            continue

        # creates new slide
        slide = prez.slides.add_slide(prez.slide_layouts[6]) # pos=6, Blank
        
        # adds title to slide
        title = slide.shapes.add_textbox(left, top, width, height)
        tf = title.text_frame
        # tf.text = img_basename
        # tf.fit_text(max_size=32) # OSError: unsupported OS, no installed fonts due to dev in WSL
        p = tf.add_paragraph()
        p.text = img_basename
        p.font.bold = True
        p.font.size = Pt(32)

        # adds image to slide
        image = slide.shapes.add_picture(img_path, left, img_top, height=img_height)
        info(f'Finished creating slide for {img_basename}')

    pptx_name = f'{dirname}.pptx'
    prez.save(pptx_name)
    info(f'Slides saved to {os.path.abspath(pptx_name)}')

    # for i, x in enumerate(prez.slide_layouts):
    #     print(i, x.name)
    #
    # 0 Title Slide
    # 1 Title and Content
    # 2 Section Header
    # 3 Two Content
    # 4 Comparison
    # 5 Title Only
    # 6 Blank
    # 7 Content with Caption
    # 8 Picture with Caption
    # 9 Title and Vertical Text
    # 10 Vertical Title and Text
        
    # slide = prez.slides.add_slide(prez.slide_layouts[8])
    # for x in slide.placeholders:
    #     print(x.name)


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description=DESCRIPTION, epilog=EPILOG)
    parser.add_argument('dir',
                        help='relative path to directory of images')
    parser.add_argument('-v', '--verbose',
                    action='store_true',
                    help='Set logging level to DEBUG')
    args = parser.parse_args()

    if args.verbose:
        l.setLevel(logging.DEBUG)

    main()